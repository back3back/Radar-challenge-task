"""
将雷达挑战任务 3 的演示文稿转换为 PowerPoint (.pptx) 格式
使用 python-pptx 库生成 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.33)  # 16:9 比例
prs.slide_height = Inches(7.5)

# 定义颜色
COLOR_BG = RGBColor(26, 26, 46)      # #1a1a2e
COLOR_ACCENT = RGBColor(233, 69, 96)  # #e94560
COLOR_CYAN = RGBColor(0, 217, 255)    # #00d9ff
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_GRAY = RGBColor(170, 170, 170)

def add_title_slide(prs, title, subtitle, info_lines=None):
    """添加标题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式

    # 背景
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_BG
    shape.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # 附加信息
    if info_lines:
        info_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(11.33), Inches(1.5))
        tf = info_box.text_frame
        for i, line in enumerate(info_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.font.color.rgb = COLOR_GRAY
            p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, content_items, two_column=False):
    """添加内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    # 标题下划线
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.1), Inches(3), Inches(0.1)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    # 内容
    if two_column:
        # 左栏
        left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(5.5))
        tf = left_box.text_frame
        items = content_items.get('left', [])
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if isinstance(item, dict):
                p.text = item.get('title', '')
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = COLOR_CYAN
            else:
                p.text = '  • ' + str(item)
                p.font.size = Pt(16)
                p.font.color.rgb = COLOR_WHITE
            p.space_after = Pt(8)

        # 右栏
        right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(6), Inches(5.5))
        tf = right_box.text_frame
        items = content_items.get('right', [])
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if isinstance(item, dict):
                p.text = item.get('title', '')
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = COLOR_CYAN
            else:
                p.text = '  • ' + str(item)
                p.font.size = Pt(16)
                p.font.color.rgb = COLOR_WHITE
            p.space_after = Pt(8)
    else:
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.5))
        tf = content_box.text_frame
        for i, item in enumerate(content_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            if isinstance(item, dict):
                p.text = item.get('title', '')
                p.font.size = Pt(20)
                p.font.bold = True
                p.font.color.rgb = COLOR_CYAN
            else:
                p.text = '  • ' + str(item)
                p.font.size = Pt(16)
                p.font.color.rgb = COLOR_WHITE
            p.space_after = Pt(8)

    return slide

def add_table_slide(prs, title, headers, data):
    """添加表格幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    # 表格
    rows = len(data) + 1
    cols = len(headers)
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(12.33)
    height = Inches(0.8)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # 设置列宽
    for i in range(cols):
        table.columns[i].width = width // cols

    # 填充表头
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(47, 52, 67)
        cell.text_frame.paragraphs[0].font.color.rgb = COLOR_ACCENT
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)

    # 填充数据
    for row_idx, row_data in enumerate(data):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(30, 30, 50)

    return slide

# ========== 创建幻灯片 ==========

# 第 1 页：标题
add_title_slide(
    prs,
    "雷达挑战任务 3",
    "阵列信号处理与测角算法仿真",
    ["汇报人：雷达项目组", "日期：2026 年 4 月"]
)

# 第 2 页：任务概述
content = {
    'left': [
        {'title': '任务 3-1: 单目标测角'},
        '仿真 1 发 4 收阵列 (d = λ/2)',
        '产生回波信号',
        'FFT 测角法仿真',
        '相位法测角仿真',
        '分析对比作用距离'
    ],
    'right': [
        {'title': '任务 3-2: 双目标测角'},
        '仿真 2 个点目标的角度',
        '分析参数对分辨力的影响',
        '阵元个数影响分析',
        '阵元间距影响分析',
        '提高分辨率方法研究'
    ]
}
add_content_slide(prs, "任务概述", content, two_column=True)

# 第 3 页：系统参数
headers = ["参数", "符号", "数值", "单位"]
data = [
    ["载波频率", "f₀", "10", "GHz"],
    ["带宽", "BW", "100", "MHz"],
    ["波长", "λ", "30", "mm"],
    ["阵元数量", "N", "4", "-"],
    ["阵元间距", "d", "λ/2", "-"],
    ["目标距离", "R", "10", "km"],
    ["目标角度", "θ", "30", "deg"]
]
add_table_slide(prs, "系统参数设计", headers, data)

# 第 4 页：信号模型
add_content_slide(prs, "信号模型", [
    {'title': '回波信号模型'},
    'sₙ(t) = s₀(t) · exp(j·2π/λ · n·d·sin(θ))',
    '',
    {'title': '相邻阵元相位差'},
    'Δφ = 2π/λ · d · sin(θ)',
    '',
    {'title': '关键参数'},
    'n = 0,1,2,...,N-1 阵元索引',
    'd 为阵元间距',
    'θ 为入射角（相对于阵列法线）'
])

# 第 5 页：测角算法原理
content = {
    'left': [
        {'title': 'FFT 测角法'},
        '空间 FFT 变换',
        '谱峰搜索定角度',
        '计算简单，实时性好',
        '可处理多目标',
        '分辨率受瑞利限限制'
    ],
    'right': [
        {'title': '相位法测角'},
        '利用相邻阵元相位差',
        '直接计算目标角度',
        '单目标精度高',
        '对噪声敏感',
        '仅适用于单目标'
    ]
}
add_content_slide(prs, "测角算法原理", content, two_column=True)

# 第 6 页：3-1 结果
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "任务 3-1 仿真结果"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT

# 添加图片占位符
placeholder = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(11.33), Inches(5.5)
)
placeholder.fill.solid()
placeholder.fill.fore_color.rgb = RGBColor(30, 30, 50)
placeholder.line.color.rgb = COLOR_CYAN
placeholder.line.width = Pt(2)

# 添加说明文字
info_box = slide.shapes.add_textbox(Inches(1), Inches(7.2), Inches(11.33), Inches(0.5))
tf = info_box.text_frame
p = tf.paragraphs[0]
p.text = "请查看 docs/p3-1-results.png"
p.font.size = Pt(14)
p.font.color.rgb = COLOR_GRAY
p.alignment = PP_ALIGN.CENTER

# 第 7 页：3-1 数据总结
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "任务 3-1 数据总结"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT

# 测角精度表格
table_data = [
    ["FFT 测角", "30.39°", "0.39°"],
    ["相位法", "30.00°", "0.00°"]
]
table = slide.shapes.add_table(3, 3, Inches(0.5), Inches(1.5), Inches(6), Inches(2)).table
headers = ["方法", "估计角度", "误差"]
for i, h in enumerate(headers):
    table.cell(0, i).text = h
    table.cell(0, i).fill.solid()
    table.cell(0, i).fill.fore_color.rgb = COLOR_ACCENT
    table.cell(0, i).text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE

for i, row in enumerate(table_data):
    for j, cell in enumerate(row):
        table.cell(i+1, j).text = cell
        table.cell(i+1, j).text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE

# 结论文字
conclusions = [
    "相邻阵元相位差：90°（与理论一致）",
    "FFT 测角：SNR ≥ 15 dB",
    "相位法：SNR ≥ 20 dB（对噪声更敏感）"
]
box = slide.shapes.add_textbox(Inches(7), Inches(1.5), Inches(6), Inches(4))
tf = box.text_frame
for i, c in enumerate(conclusions):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = "✓ " + c
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(74, 222, 128)

# 第 8 页：3-2 结果
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "任务 3-2 仿真结果"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT

placeholder = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(11.33), Inches(5.5)
)
placeholder.fill.solid()
placeholder.fill.fore_color.rgb = RGBColor(30, 30, 50)
placeholder.line.color.rgb = COLOR_CYAN
placeholder.line.width = Pt(2)

info_box = slide.shapes.add_textbox(Inches(1), Inches(7.2), Inches(11.33), Inches(0.5))
tf = info_box.text_frame
p = tf.paragraphs[0]
p.text = "请查看 docs/p3-2-results.png"
p.font.size = Pt(14)
p.font.color.rgb = COLOR_GRAY
p.alignment = PP_ALIGN.CENTER

# 第 9 页：角度分辨力分析
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "角度分辨力分析"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT

# 瑞利判据公式
formula_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12.33), Inches(0.8))
tf = formula_box.text_frame
p = tf.paragraphs[0]
p.text = "Δθ ≈ λ / L = λ / [(N-1)·d]"
p.font.size = Pt(24)
p.font.color.rgb = COLOR_CYAN

# 分辨力表格
headers = ["阵元数", "d=0.25λ", "d=0.50λ", "d=0.75λ", "d=1.00λ"]
data = [
    ["4", "76.39°", "38.20°", "25.46°", "19.10°"],
    ["8", "32.74°", "16.37°", "10.91°", "8.19°"],
    ["16", "15.28°", "7.64°", "5.09°", "3.82° ✓"]
]
table = slide.shapes.add_table(4, 5, Inches(0.5), Inches(2.3), Inches(12.33), Inches(2.5)).table

for i, h in enumerate(headers):
    cell = table.cell(0, i)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = COLOR_ACCENT
    cell.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
    cell.text_frame.paragraphs[0].font.size = Pt(12)

for row_idx, row in enumerate(data):
    for col_idx, val in enumerate(row):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        cell.text_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
        cell.text_frame.paragraphs[0].font.size = Pt(11)

# 结论文字
warning_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(12.33), Inches(0.8))
tf = warning_box.text_frame
p = tf.paragraphs[0]
p.text = "⚠️ 4 阵元阵列无法分辨 5°间隔的双目标，需要至少 16 阵元"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(251, 191, 36)

# 第 10 页：提高分辨率的方法
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "提高分辨率的方法"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT

methods = [
    ("MUSIC 算法", ["子空间分解", "突破瑞利限", "需已知信源数", "计算量大"]),
    ("ESPRIT 算法", ["旋转不变性", "无需谱搜索", "计算量小", "需阵列校准"]),
    ("压缩感知", ["角度域稀疏性", "少快拍适用", "L1 最小化", "算法复杂"]),
    ("最大似然", ["统计最优", "精度最高", "多维搜索", "复杂度极高"])
]

for i, (title, items) in enumerate(methods):
    left = Inches(0.5 + (i % 2) * 6.8)
    top = Inches(1.5 + (i // 2) * 2.8)

    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(6.5), Inches(2.5)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(30, 30, 50)
    box.line.color.rgb = COLOR_CYAN
    box.line.width = Pt(2)

    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    for j, item in enumerate(items):
        p = tf.add_paragraph()
        p.text = "  • " + item
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_WHITE

# 第 11 页：代码结构
add_content_slide(prs, "代码结构", [
    "src/radar/",
    "  ├── p3-1.py      # 单目标测角",
    "  ├── p3-2.py      # 双目标测角",
    "  └── __init__.py",
    "",
    "tests/",
    "  └── test_radar.py  # 单元测试",
    "",
    "核心函数:",
    "  • generate_echo_signal()",
    "  • fft_angle_estimation()",
    "  • phase_angle_estimation()",
    "  • music_angle_estimation()",
    "  • esprit_angle_estimation()"
])

# 第 12 页：测试结果
slide = prs.slides.add_slide(prs.slide_layouts[6])
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.33), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "测试结果"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT

# 测试结果统计
result_box = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(5), Inches(2.5)
)
result_box.fill.solid()
result_box.fill.fore_color.rgb = RGBColor(30, 30, 50)
result_box.line.color.rgb = RGBColor(74, 222, 128)
result_box.line.width = Pt(3)

tf = result_box.text_frame
p = tf.paragraphs[0]
p.text = "19 passed in 1.33s"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(74, 222, 128)
p.alignment = PP_ALIGN.CENTER

# 测试列表
tests = [
    "✓ 回波信号生成测试",
    "✓ FFT 测角精度测试",
    "✓ 相位法测角测试",
    "✓ MUSIC 算法测试",
    "✓ ESPRIT 算法测试",
    "✓ 角分辨力计算测试",
    "✓ 集成测试"
]
box = slide.shapes.add_textbox(Inches(6.5), Inches(1.5), Inches(6), Inches(4))
tf = box.text_frame
for i, t in enumerate(tests):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = t
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(74, 222, 128)

# 第 13 页：结论
content = {
    'left': [
        {'title': '任务 3-1 结论'},
        '✓ FFT 和相位法均能准确测角',
        '相位法单目标精度更高',
        'FFT 法噪声抑制更好',
        '仿真结果与理论一致'
    ],
    'right': [
        {'title': '任务 3-2 结论'},
        '4 阵元分辨力约 38.2°',
        '无法分辨 5°间隔目标',
        '⚠️ 需 16 阵元才能分辨 5°间隔',
        '高分辨率算法可突破限制'
    ]
}
add_content_slide(prs, "结论", content, two_column=True)

# 第 14 页：结束
add_title_slide(
    prs,
    "感谢聆听",
    "欢迎提问",
    ["详细报告：挑战任务 3.md", "仿真代码：src/radar/p3-1.py, p3-2.py"]
)

# 保存 PPTX
output_path = os.path.join(os.path.dirname(__file__), 'presentation.pptx')
prs.save(output_path)
print(f"PPTX 已保存至：{output_path}")
